#!/usr/bin/env python3
import argparse
import json
import math
from pathlib import Path
from collections import defaultdict

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


DATASET_DISPLAY = {
    "flashrag_nq": "NQ",
    "flashrag_triviaqa": "TriviaQA",
    "flashrag_popqa": "PopQA",
    "flashrag_hotpotqa": "HotpotQA",
    "flashrag_2wikimultihopqa": "2wiki",
    "flashrag_musique": "Musique",
    "flashrag_bamboogle": "Bamboogle",
}

DEFAULT_BASE = Path("/root/code/SeedMitigating/output/eval_flashrag500_7models")
DEFAULT_ROOT = DEFAULT_BASE / "response_level"
DEFAULT_OUT_DIR = DEFAULT_BASE / "report"
DEFAULT_OUT_PPT = DEFAULT_OUT_DIR / "eval_flashrag500_7models_report.pptx"
DEFAULT_OUT_SUMMARY = DEFAULT_OUT_DIR / "analysis_summary.txt"
METRIC_DIRECTIONS = {
    "SNR Gain": "up",
    "Conf AUC": "up",
    "Abs Acc": "up",
    "smECE": "down",
    "Brier": "down",
    "NLL": "down",
    "Pred Acc": "up",
}


def to_num(x, default=0.0):
    try:
        return float(x)
    except Exception:
        return default


def parse_args():
    parser = argparse.ArgumentParser(
        description="Export PPT and summary from a response_level results directory."
    )
    parser.add_argument(
        "--root",
        type=Path,
        default=DEFAULT_ROOT,
        help="Input response_level directory. Default: %(default)s",
    )
    parser.add_argument(
        "--out-ppt",
        type=Path,
        default=DEFAULT_OUT_PPT,
        help="Output pptx path. Default: %(default)s",
    )
    parser.add_argument(
        "--out-summary",
        type=Path,
        default=DEFAULT_OUT_SUMMARY,
        help="Output summary txt path. Default: %(default)s",
    )
    parser.add_argument(
        "--skip-ppt",
        action="store_true",
        help="Do not export pptx.",
    )
    parser.add_argument(
        "--skip-summary",
        action="store_true",
        help="Do not export summary txt.",
    )
    return parser.parse_args()


def build_data(root: Path):
    rows = []
    image_rows = []

    for method_dir in sorted([p for p in root.iterdir() if p.is_dir()]):
        method = method_dir.name
        for dataset_dir in sorted([p for p in method_dir.iterdir() if p.is_dir()]):
            dataset = dataset_dir.name
            result_file = dataset_dir / "results.json"
            if not result_file.exists():
                continue

            with open(result_file, "r", encoding="utf-8") as f:
                data = json.load(f)

            metrics = data.get("metrics") or data.get("metrics_confidence") or {}
            row = {
                "method": method,
                "dataset": dataset,
                "SNR Gain": to_num(metrics.get("snr_gain")),
                "Conf AUC": to_num(metrics.get("confidence_auc")),
                "Abs Acc": to_num(metrics.get("abstention_accuracy")),
                "smECE": to_num(metrics.get("smoothed_ece")),
                "Brier": to_num(metrics.get("brier_score")),
                "NLL": to_num(metrics.get("nll")),
                "Pred Acc": to_num(metrics.get("predictive_accuracy")),
                "Avg Conf": to_num(metrics.get("avg_confidence")),
                "Abstention Rate": to_num(metrics.get("abstention_rate")),
                "n_samples": int(metrics.get("n_samples", 0) or 0),
                "n_answered": int(metrics.get("n_answered", 0) or 0),
                "n_abstained": int(metrics.get("n_abstained", 0) or 0),
                "results_json": str(result_file),
            }
            rows.append(row)

            for img in sorted(dataset_dir.glob("*.png")):
                image_rows.append({"method": method, "dataset": dataset, "image": str(img)})
            for img in sorted(dataset_dir.glob("*.jpg")):
                image_rows.append({"method": method, "dataset": dataset, "image": str(img)})
            for img in sorted(dataset_dir.glob("*.jpeg")):
                image_rows.append({"method": method, "dataset": dataset, "image": str(img)})

    df = pd.DataFrame(rows)
    images = pd.DataFrame(image_rows)
    return df, images


def get_method_order(df: pd.DataFrame):
    preferred = [
        "Qwen3-4B-Instruct-2507 (Base)",
        "Binary Reward RLOO",
        "Verbalized Brier RLOO",
        "Binary Reward (TriviaQA)",
        "Verbalized CE (TriviaQA)",
        "Verbalized Brier (TriviaQA)",
        "Binary Reward (NQ+HotpotQA)",
        "Verbalized CE (NQ+HotpotQA)",
        "Verbalized Brier (NQ+HotpotQA)",
    ]
    existing = df["method"].unique().tolist()
    ordered = []
    for prefix in preferred:
        ordered.extend([m for m in existing if m == prefix or m.startswith(prefix)])
    ordered += [m for m in existing if m not in ordered]
    return ordered


def get_dataset_order(df: pd.DataFrame):
    preferred = [
        "flashrag_nq",
        "flashrag_triviaqa",
        "flashrag_popqa",
        "flashrag_hotpotqa",
        "flashrag_2wikimultihopqa",
        "flashrag_musique",
        "flashrag_bamboogle",
    ]
    existing = df["dataset"].unique().tolist()
    ordered = [d for d in preferred if d in existing]
    ordered += [d for d in existing if d not in ordered]
    return ordered


def compute_tables(df: pd.DataFrame):
    metric_cols = ["SNR Gain", "Conf AUC", "Abs Acc", "smECE", "Brier", "NLL", "Pred Acc"]
    up_metrics = {"SNR Gain", "Conf AUC", "Abs Acc", "Pred Acc"}
    down_metrics = {"smECE", "Brier", "NLL"}

    method_order = get_method_order(df)
    dataset_order = get_dataset_order(df)

    # coverage
    coverage = (
        df.assign(done=1)
        .pivot_table(index="method", columns="dataset", values="done", aggfunc="max", fill_value=0)
        .reindex(index=method_order, columns=dataset_order)
    )

    # by dataset tables
    by_dataset = {}
    for ds in dataset_order:
        t = (
            df[df["dataset"] == ds]
            .set_index("method")[metric_cols]
            .reindex(method_order)
        )
        by_dataset[ds] = t

    # overall mean
    overall = (
        df.groupby("method", as_index=True)[metric_cols]
        .mean()
        .reindex(method_order)
    )

    # win count (per dataset)
    wins = defaultdict(lambda: defaultdict(int))
    for ds, t in by_dataset.items():
        for m in metric_cols:
            s = t[m].dropna()
            if s.empty:
                continue
            if m in up_metrics:
                best_val = s.max()
                winners = s[s == best_val].index.tolist()
            else:
                best_val = s.min()
                winners = s[s == best_val].index.tolist()
            for w in winners:
                wins[w][m] += 1

    win_df = pd.DataFrame(index=method_order)
    for m in metric_cols:
        win_df[m] = [wins[method].get(m, 0) for method in method_order]
    win_df["Total"] = win_df.sum(axis=1)

    return {
        "metric_cols": metric_cols,
        "up_metrics": up_metrics,
        "down_metrics": down_metrics,
        "method_order": method_order,
        "dataset_order": dataset_order,
        "coverage": coverage,
        "by_dataset": by_dataset,
        "overall": overall,
        "win_df": win_df,
    }


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle_box = slide.placeholders[1]
    subtitle_box.text = ""
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    lines = subtitle.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11 if i == 0 else 9)
        p.font.bold = False


def add_text_slide(prs: Presentation, title: str, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(0.5)
    top = Inches(1.2)
    width = prs.slide_width - Inches(1.0)
    height = prs.slide_height - Inches(1.6)
    tx = slide.shapes.add_textbox(left, top, width, height).text_frame
    tx.word_wrap = True
    for i, line in enumerate(lines):
        p = tx.paragraphs[0] if i == 0 else tx.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if i == 0 else 14)


def add_table_slide(prs: Presentation, title: str, df: pd.DataFrame, metric_cols, up_metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    rows = len(df) + 1
    cols = len(df.columns)
    left = Inches(0.3)
    top = Inches(1.25)
    width = prs.slide_width - Inches(0.6)
    height = prs.slide_height - Inches(2.0)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # set column widths
    first_col_w = int(width * 0.3)
    other_col_w = int((width - first_col_w) / max(1, cols - 1))
    for c in range(cols):
        table.columns[c].width = first_col_w if c == 0 else other_col_w

    # header
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(217, 225, 242)

    # find best values per metric
    best_map = {}
    for m in metric_cols:
        display_col = metric_display_name(m)
        if display_col in df.columns:
            s = pd.to_numeric(df[display_col], errors="coerce").dropna()
            if not s.empty:
                best_map[display_col] = s.max() if m in up_metrics else s.min()

    # data cells
    for r in range(len(df)):
        for c, col in enumerate(df.columns):
            v = df.iloc[r, c]
            cell = table.cell(r + 1, c)
            if isinstance(v, float):
                txt = f"{v:.4f}"
            else:
                txt = str(v)
            cell.text = txt
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT

            # bold best
            if col in best_map:
                vv = pd.to_numeric(pd.Series([v]), errors="coerce").iloc[0]
                if pd.notna(vv) and abs(vv - best_map[col]) < 1e-12:
                    p.font.bold = True


def add_image_slide(prs: Presentation, title: str, image_path: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    left = Inches(0.3)
    top = Inches(1.0)
    max_w = prs.slide_width - Inches(0.6)
    max_h = prs.slide_height - Inches(1.3)

    pic = slide.shapes.add_picture(str(image_path), left, top)
    # scale to fit
    scale = min(1.0, max_w / pic.width, max_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int((prs.slide_width - pic.width) / 2)
    pic.top = int(top + (max_h - pic.height) / 2)


def simplify_method_label(method: str) -> str:
    if method == "Qwen3-4B-Instruct-2507 (Base)":
        return "Base"
    if method.startswith("Binary Reward RLOO"):
        return method.replace("Binary Reward RLOO", "Binary Reward RLOO").strip()
    if method.startswith("Verbalized Brier RLOO"):
        return method.replace("Verbalized Brier RLOO", "Verbalized Brier RLOO").strip()
    if method.startswith("Binary Reward"):
        return "Binary Reward"
    if method.startswith("Verbalized Brier"):
        return "Verbalized Brier"
    if method.startswith("Verbalized CE"):
        return "Verbalized CE"
    return method


def dataset_label(dataset: str) -> str:
    return DATASET_DISPLAY.get(dataset, dataset.replace("flashrag_", ""))


def metric_display_name(metric: str) -> str:
    direction = METRIC_DIRECTIONS.get(metric)
    if direction == "up":
        return f"{metric} (↑)"
    if direction == "down":
        return f"{metric} (↓)"
    return metric


def infer_train_set(method: str) -> str:
    if method == "Qwen3-4B-Instruct-2507 (Base)":
        return "-"
    if "triviaqa" in method.lower():
        return "TriviaQA"
    if "hotpotqa" in method.lower():
        return "HotpotQA"
    return "TriviaQA"


def infer_data_size(method: str) -> str:
    if method == "Qwen3-4B-Instruct-2507 (Base)":
        return "-"
    if "RLOO" in method:
        return "10k / 2k"
    return "-"


def method_metadata(method: str) -> dict:
    return {
        "Method": simplify_method_label(method),
        "Train Set": infer_train_set(method),
        "Data Size": infer_data_size(method),
    }


def infer_test_size(df: pd.DataFrame, dataset: str) -> str:
    subset = df[df["dataset"] == dataset]
    if subset.empty or "n_samples" not in subset.columns:
        return "-"
    vals = [int(v) for v in subset["n_samples"].dropna().tolist()]
    if not vals:
        return "-"
    uniq = sorted(set(vals))
    if len(uniq) == 1:
        return str(uniq[0])
    return "/".join(str(v) for v in uniq)


def add_fitted_image(slide, image_path: Path, box_left, box_top, box_w, box_h):
    pic = slide.shapes.add_picture(str(image_path), box_left, box_top)
    scale = min(box_w / pic.width, box_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(box_left + (box_w - pic.width) / 2)
    pic.top = int(box_top + (box_h - pic.height) / 2)


def add_calibration_grid(slide, methods, dataset: str, image_lookup: dict, content_left, content_top, content_w, content_h):
    methods = list(methods)
    if not methods:
        return

    cols = 2
    rows = max(1, math.ceil(len(methods) / cols))
    gap_x = Inches(0.05)
    gap_y = Inches(0.05)
    cell_w = int((content_w - gap_x * (cols - 1)) / cols)
    cell_h = int((content_h - gap_y * (rows - 1)) / rows)
    caption_h = Inches(0.16)
    img_h = int(cell_h - caption_h)

    for idx, method in enumerate(methods):
        r = idx // cols
        c = idx % cols
        cell_left = int(content_left + c * (cell_w + gap_x))
        cell_top = int(content_top + r * (cell_h + gap_y))

        img_path = image_lookup.get((method, dataset))
        if img_path and img_path.exists():
            add_fitted_image(slide, img_path, cell_left, cell_top, cell_w, img_h)
        else:
            tx = slide.shapes.add_textbox(cell_left, cell_top, cell_w, img_h).text_frame
            tx.text = "missing confidence_calibration.png"
            tx.paragraphs[0].font.size = Pt(12)
            tx.paragraphs[0].alignment = PP_ALIGN.CENTER

        cap = slide.shapes.add_textbox(cell_left, int(cell_top + img_h), cell_w, caption_h).text_frame
        meta = method_metadata(method)
        cap.text = meta["Method"]
        cp = cap.paragraphs[0]
        cp.font.size = Pt(8)
        cp.alignment = PP_ALIGN.CENTER


def add_dataset_combo_slide(prs: Presentation, title: str, df: pd.DataFrame, methods, dataset: str, test_size: str, image_lookup: dict, metric_cols, up_metrics):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.12), Inches(0.04), prs.slide_width - Inches(0.24), Inches(0.22))
    tf = title_box.text_frame
    tf.word_wrap = False
    tf.text = f"{title} | Test Size: {test_size}"
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    meta_lines = []
    for method in methods:
        meta = method_metadata(method)
        meta_lines.append(f"{meta['Method']} | Train: {meta['Train Set']} | Data: {meta['Data Size']}")
    meta_box = slide.shapes.add_textbox(Inches(0.18), Inches(0.27), prs.slide_width - Inches(0.36), Inches(0.32))
    meta_tf = meta_box.text_frame
    meta_tf.word_wrap = True
    for idx, line in enumerate(meta_lines):
        para = meta_tf.paragraphs[0] if idx == 0 else meta_tf.add_paragraph()
        para.text = line
        para.font.size = Pt(9)
        para.font.bold = True
        para.alignment = PP_ALIGN.CENTER

    table_left = Inches(0.18)
    table_top = Inches(0.65)
    table_w = prs.slide_width - Inches(0.36)
    table_h = Inches(1.45)
    table_df = df.copy()
    table_df.columns = [metric_display_name(c) if c in METRIC_DIRECTIONS else c for c in table_df.columns]
    add_table_on_slide(slide, table_df, table_left, table_top, table_w, table_h, metric_cols, up_metrics)

    content_left = Inches(0.12)
    content_top = Inches(2.25)
    content_right_margin = Inches(0.12)
    content_bottom_margin = Inches(0.10)
    content_w = prs.slide_width - content_left - content_right_margin
    content_h = prs.slide_height - content_top - content_bottom_margin
    add_calibration_grid(slide, methods, dataset, image_lookup, content_left, content_top, content_w, content_h)


def add_table_on_slide(slide, df: pd.DataFrame, left, top, width, height, metric_cols, up_metrics):
    rows = len(df) + 1
    cols = len(df.columns)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    first_col_w = int(width * 0.18)
    second_col_w = int(width * 0.12)
    third_col_w = int(width * 0.12)
    rest_w = int((width - first_col_w - second_col_w - third_col_w) / max(1, cols - 3))
    for c in range(cols):
        if c == 0:
            table.columns[c].width = first_col_w
        elif c == 1:
            table.columns[c].width = second_col_w
        elif c == 2:
            table.columns[c].width = third_col_w
        else:
            table.columns[c].width = rest_w

    best_map = {}
    for metric in metric_cols:
        label = metric_display_name(metric)
        if label in df.columns:
            s = pd.to_numeric(df[label], errors="coerce").dropna()
            if not s.empty:
                best_map[label] = s.max() if metric in up_metrics else s.min()

    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(9)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(217, 225, 242)

    for r in range(len(df)):
        for c, col in enumerate(df.columns):
            v = df.iloc[r, c]
            cell = table.cell(r + 1, c)
            cell.text = f"{v:.4f}" if isinstance(v, float) else str(v)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if col in best_map:
                vv = pd.to_numeric(pd.Series([v]), errors="coerce").iloc[0]
                if pd.notna(vv) and abs(vv - best_map[col]) < 1e-12:
                    p.font.bold = True


def chunked(items, size):
    for i in range(0, len(items), size):
        yield items[i:i + size]


def write_ppt(output_ppt: Path, df: pd.DataFrame, images: pd.DataFrame, tables: dict, root: Path):
    prs = Presentation()
    method_order = tables["method_order"]
    dataset_order = tables["dataset_order"]
    metric_cols = tables["metric_cols"]
    up_metrics = tables["up_metrics"]

    subtitle_lines = [f"Methods: {len(method_order)} | Datasets: {len(dataset_order)} | Rows: {len(df)}"]
    for method in method_order:
        meta = method_metadata(method)
        subtitle_lines.append(f"{meta['Method']} | Train: {meta['Train Set']} | Data: {meta['Data Size']}")

    add_title_slide(
        prs,
        "FlashRAG500 Evaluation Report",
        "\n".join(subtitle_lines),
    )

    overall_df = tables["overall"].reset_index().rename(columns={"method": "Method"})
    overall_df.insert(1, "Train Set", overall_df["Method"].map(lambda x: method_metadata(x)["Train Set"]))
    overall_df.insert(2, "Data Size", overall_df["Method"].map(lambda x: method_metadata(x)["Data Size"]))
    overall_df["Method"] = overall_df["Method"].map(simplify_method_label)
    overall_df.columns = [metric_display_name(c) if c in METRIC_DIRECTIONS else c for c in overall_df.columns]
    add_table_slide(prs, "Overall Mean Metrics", overall_df, metric_cols, up_metrics)

    win_df = tables["win_df"].reset_index().rename(columns={"index": "Method"})
    win_df.insert(1, "Train Set", win_df["Method"].map(lambda x: method_metadata(x)["Train Set"]))
    win_df.insert(2, "Data Size", win_df["Method"].map(lambda x: method_metadata(x)["Data Size"]))
    win_df["Method"] = win_df["Method"].map(simplify_method_label)
    add_table_slide(prs, "Win Count Across Datasets", win_df, metric_cols, up_metrics)

    image_lookup = {}
    if not images.empty:
        images = images.sort_values(["dataset", "method", "image"]).reset_index(drop=True)
        for _, row in images.iterrows():
            img = Path(row["image"])
            if img.name not in {"confidence_calibration.png", "conf_calibration.png"}:
                continue
            image_lookup[(row["method"], row["dataset"])] = img

    for ds in dataset_order:
        method_groups = list(chunked(method_order, 4))
        for page_idx, methods in enumerate(method_groups, start=1):
            suffix = f" ({page_idx}/{len(method_groups)})" if len(method_groups) > 1 else ""
            title = f"Test: {dataset_label(ds)}{suffix}"
            test_size = infer_test_size(df, ds)
            dataset_df = tables["by_dataset"][ds].loc[list(methods)].reset_index().rename(columns={"method": "Method"})
            dataset_df["Method"] = dataset_df["Method"].map(simplify_method_label)
            add_dataset_combo_slide(prs, title, dataset_df, methods, ds, test_size, image_lookup, metric_cols, up_metrics)

    prs.save(output_ppt)


def write_summary(summary_path: Path, df: pd.DataFrame, tables: dict):
    metric_cols = tables["metric_cols"]
    up_metrics = tables["up_metrics"]
    overall = tables["overall"]

    lines = []
    lines.append("FlashRAG500 evaluation summary")
    lines.append(f"Rows: {len(df)} | Methods: {df['method'].nunique()} | Datasets: {df['dataset'].nunique()}")
    lines.append("")
    lines.append("Overall best by metric:")
    for m in metric_cols:
        s = overall[m].dropna()
        if s.empty:
            continue
        best_val = s.max() if m in up_metrics else s.min()
        winners = s[s == best_val].index.tolist()
        lines.append(f"- {m}: {', '.join(winners)} ({best_val:.4f})")

    win_df = tables["win_df"]
    total_best = win_df["Total"].max()
    total_winners = win_df[win_df["Total"] == total_best].index.tolist()
    lines.append("")
    lines.append(f"Best total wins: {', '.join(total_winners)} ({int(total_best)})")

    summary_path.write_text("\n".join(lines), encoding="utf-8")


def main():
    args = parse_args()
    root = args.root.resolve()
    output_ppt = args.out_ppt.resolve()
    output_summary = args.out_summary.resolve()
    output_ppt.parent.mkdir(parents=True, exist_ok=True)
    output_summary.parent.mkdir(parents=True, exist_ok=True)

    df, images = build_data(root)
    if df.empty:
        raise SystemExit(f"No results found under {root}")

    tables = compute_tables(df)

    if not args.skip_ppt:
        write_ppt(output_ppt, df, images, tables, root)
    if not args.skip_summary:
        write_summary(output_summary, df, tables)

    if not args.skip_ppt:
        print(f"[OK] PPTX: {output_ppt}")
    if not args.skip_summary:
        print(f"[OK] SUMMARY: {output_summary}")


if __name__ == "__main__":
    main()
